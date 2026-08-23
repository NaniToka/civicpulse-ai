from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()

# Brand colors based on CivicPulse UI
COLOR_PRIMARY = RGBColor(79, 70, 229)   # Indigo 600
COLOR_DARK = RGBColor(15, 23, 42)       # Slate 900
COLOR_TEXT = RGBColor(51, 65, 85)       # Slate 700

def add_title_slide(prs, title_text, subtitle_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = subtitle_text
    subtitle.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

def add_content_slide(prs, title_text, bullet_points):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body_shape = slide.shapes.placeholders[1]
    
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = COLOR_DARK
    title.text_frame.paragraphs[0].font.bold = True
    
    tf = body_shape.text_frame
    tf.text = bullet_points[0]
    tf.paragraphs[0].font.color.rgb = COLOR_TEXT
    tf.paragraphs[0].font.size = Pt(24)
    
    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.font.color.rgb = COLOR_TEXT
        p.font.size = Pt(24)
        p.level = 0

# Slide 1: Title
add_title_slide(
    prs, 
    "CivicPulse AI", 
    "Citizen Demand Intelligence & Infrastructure Prioritization\nTransforming Civic Feedback into Actionable Budget Plans"
)

# Slide 2: The Problem
add_content_slide(prs, "The Disconnect: Budgets vs. Real Needs", [
    "Governments spend billions on infrastructure, but allocations often rely on outdated census data.",
    "Citizens lack a centralized, accessible platform in their native languages to report local issues (e.g., broken water pipes, missing hospitals).",
    "Result: Highly inefficient budget spending and ignored communities."
])

# Slide 3: The Solution
add_content_slide(prs, "Introducing CivicPulse AI", [
    "A real-time, AI-powered command center that listens to citizens directly.",
    "Synthesizes thousands of unstructured complaints into actionable, prioritized infrastructure projects.",
    "Bridges the gap between raw citizen demand and government capital allocation."
])

# Slide 4: Key Feature 1
add_content_slide(prs, "1. Multilingual Citizen Voice Studio", [
    "Breaks the language barrier for rural and urban populations alike.",
    "Citizens can submit text or voice feedback in 7+ native languages (Telugu, Hindi, Marathi, Bengali, Tamil, etc.).",
    "AI automatically translates the input, analyzes sentiment urgency, and categorizes the problem."
])

# Slide 5: Key Feature 2
add_content_slide(prs, "2. Demand Intelligence & Hotspot Mapping", [
    "Maps infrastructure deficits geographically in real-time.",
    "Groups incoming complaints to identify 'Problem Hotspots' across monitored districts.",
    "Generates a 'Facility Shortfall Matrix' to show exactly where existing infrastructure is failing."
])

# Slide 6: Key Feature 3
add_content_slide(prs, "3. AI Budget & Investment Simulator", [
    "Optimizes capital deployment for policymakers.",
    "Input a proposed budget (e.g., $15M) and simulate the exact impact of building new facilities.",
    "Instantly calculates ROI metrics: Estimated citizens benefited, cost per citizen, and deficit reduction percentages."
])

# Slide 7: Key Feature 4
add_content_slide(prs, "4. Transparent Evidence Trails", [
    "Explainable AI builds trust between government officials and the public.",
    "Why was a hospital recommended in Kanpur instead of a transit hub in Pune?",
    "A 6-step recommendation evidence chain provides a transparent, machine-readable proof trail from citizen complaint to final policy recommendation."
])

# Slide 8: Technical Architecture
add_content_slide(prs, "How It Works Under the Hood", [
    "Frontend: React, TypeScript, Tailwind CSS, Vite (Optimized for speed and responsiveness).",
    "Backend: Python, FastAPI (Scalable and high-performance API routing).",
    "AI Engine: LLMs for translation, sentiment extraction, and spatial-demand synthesis.",
    "Data Layer: Real-time feedback aggregation and simulated census/infrastructure tracking."
])

# Slide 9: Impact & Roadmap
add_content_slide(prs, "Impact & Future Roadmap", [
    "Impact: Faster response times, data-driven budgeting, and significantly higher citizen satisfaction and trust.",
    "Future: Integration with IoT sensors (e.g., water quality monitors) for automated anomaly detection.",
    "Future: Expansion to 50+ languages and direct API integration with city contractor dispatch systems."
])

# Slide 10: Conclusion
add_title_slide(
    prs, 
    "Thank You", 
    "Empowering citizens and enabling data-driven governance.\n\nOpen for Q&A and Live Demo."
)

# Save the presentation
prs.save('/Users/tokanani/civicpulse-ai/CivicPulse_AI_Presentation.pptx')
print("Successfully generated presentation!")
